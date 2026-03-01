import os
import json
import re
import wave
from google import genai
from google.genai import types
from src.agents.states import GraphState
from src.tools.database import get_vectorstore
from langchain_text_splitters import RecursiveCharacterTextSplitter
from dotenv import load_dotenv

load_dotenv()

client = genai.Client(api_key=os.getenv("GEMINI_API_KEY"))

def parse_json_safe(text: str):
    
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        pass

    #Strip markdown code fences 
    cleaned = re.sub(r'^```(?:json)?\s*', '', text.strip(), flags=re.IGNORECASE)
    cleaned = re.sub(r'\s*```$', '', cleaned.strip())
    try:
        return json.loads(cleaned)
    except json.JSONDecodeError:
        pass

    #Extract first JSON array or object using regex
    match = re.search(r'(\[.*\]|\{.*\})', cleaned, re.DOTALL)
    if match:
        try:
            return json.loads(match.group(1))
        except json.JSONDecodeError:
            pass

    #Remove trailing commas before } or ] (common Gemini mistake)
    fixed = re.sub(r',\s*([}\]])', r'\1', cleaned)
    return json.loads(fixed)

#node function

# Chunks the PDF text and uploads it to Pinecone for RAG

def ingest_to_pinecone_node(state: GraphState):
    print("---NODE: INDEXING TO PINECONE---")
    text = state["pdf_text"]
    session_id = state.get("session_id", "default")

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
    chunks = text_splitter.split_text(text)

    vectorstore = get_vectorstore()
    metadatas = [{"session_id": session_id} for _ in chunks]
    vectorstore.add_texts(chunks, metadatas=metadatas)

    return {"next_step": "summarize"}

# """Uses Gemini to create a concise summary of the PDF."""

def summarize_pdf_node(state: GraphState):
    print("---NODE: SUMMARIZING---")
    pdf_text = state["pdf_text"]
    
    prompt = f"Summarize the following document in concise manner in a professional way. Keep main ponits and key words blend them together and provide proper summary\n\n{pdf_text}"
    response = client.models.generate_content(model="gemini-2.5-flash", contents=prompt)
    summary = response.candidates[0].content.parts[0].text
    
    return {"summary_text": summary}


# Converts the summary text into speech using Gemini TTS.
def generate_tts_node(state: GraphState):
    
    print("---NODE: GENERATING VOICE---")
    # Always ensure summary is present, generate if missing
    summary = state.get("summary_text")
    if not summary:
        pdf_text = state["pdf_text"]
        prompt = (
            "Summarize the following document in concise manner in a professional way. "
            "Keep main points and key words, blend them together and provide proper summary\n\n"
            f"{pdf_text}"
        )
        response = client.models.generate_content(model="gemini-2.5-flash", contents=prompt)
        summary = response.candidates[0].content.parts[0].text
        # Optionally update state with summary for downstream nodes
        state["summary_text"] = summary

    config = types.GenerateContentConfig(
        response_modalities=["AUDIO"],
        speech_config=types.SpeechConfig(
            voice_config=types.VoiceConfig(
                prebuilt_voice_config=types.PrebuiltVoiceConfig(voice_name="Algenib")
            )
        )
    )

    response = client.models.generate_content(
        model="gemini-2.5-pro-preview-tts",
        contents=summary,
        config=config
    )

    audio_data = response.candidates[0].content.parts[0].inline_data.data

    # Save audio file to assets folder as WAV
    os.makedirs("assets", exist_ok=True)
    audio_path = "assets/summary_audio.wav"
    with wave.open(audio_path, "wb") as wav_file:
        wav_file.setnchannels(1)
        wav_file.setsampwidth(2)   # 16-bit PCM
        wav_file.setframerate(24000)
        wav_file.writeframes(audio_data)

    return {"audio_path": audio_path, "summary_text": summary}



def generate_script_node(state: GraphState):
    print("---NODE : CREATING SCRIPT---")
    pdf_text=state["pdf_text"]
    prompt = f"""
    Transform the following PDF content into an engaging, CASUAL , 2-PERSON PODCAST STYLE SCRIPT.
    it should be a  audio clip, but should have 2 speaker's speech, one host and another expert
    Generate the script in this format as shown
    host: "host introduces ."
    Expert :" starts conversing aboy the topic in  casual way just like in podcast 2 people onverse"
    Host: "Again host interrupts the expert and speaks or asks something related to the pdf text"
    It should be an egaging like one will speak then another , then again the first person will speak.
    eACH person script will be 
    Speaker 1 (Host): Curious and energetic.
    Speaker 2 (Expert): Clear and insightful.
    
    Output ONLY a JSON list of objects with "speaker" and "text" keys.
    THE WHOLE SCRIPT GENERATED  SHOULD BE WITHIN 100-150 WORDS
    Content: {pdf_text}
    """
    response=client.models.generate_content(model="gemini-2.5-flash", 
                                            contents=prompt,
                                            config={"response_mime_type": "application/json"}
                                            )
    return{"script": response.text}


def generate_tts_script_node(state: GraphState):
    script = parse_json_safe(state["script"])
    
    # Build a plain-text dialogue prompt with speaker labels matching the config
    dialogue_prompt = "\n".join(
        f"{line['speaker']}: {line['text']}" for line in script
    )
    
    config = types.GenerateContentConfig(
        response_modalities=["AUDIO"],
        speech_config=types.SpeechConfig(
            multi_speaker_voice_config=types.MultiSpeakerVoiceConfig(
                speaker_voice_configs=[
                    types.SpeakerVoiceConfig(
                        speaker="Host",
                        voice_config=types.VoiceConfig(
                            prebuilt_voice_config=types.PrebuiltVoiceConfig(voice_name="Algenib")
                        )
                    ),
                    types.SpeakerVoiceConfig(
                        speaker="Expert",
                        voice_config=types.VoiceConfig(
                            prebuilt_voice_config=types.PrebuiltVoiceConfig(voice_name="Aoede")
                        )
                    ),
                ]
            )
        )
    )
    
    response = client.models.generate_content(
        model="gemini-2.5-pro-preview-tts",
        contents=dialogue_prompt,
        config=config
    )
    
    audio_data = response.candidates[0].content.parts[0].inline_data.data

    os.makedirs("assets", exist_ok=True)
    audio_path = "assets/podcast_summary.wav"
    with wave.open(audio_path, "wb") as wav_file:
        wav_file.setnchannels(1)
        wav_file.setsampwidth(2)   # 16-bit PCM
        wav_file.setframerate(24000)
        wav_file.writeframes(audio_data)
        
    return {"audio_path": audio_path}


#  """Generates 3 follow-up questions based on the PDF content."""
def generate_suggestions_node(state: GraphState):
   
    print("---NODE: GENERATING SUGGESTIONS---")
    pdf_text = state["pdf_text"]
    
    prompt = f"Based on this text, suggest 3 short, interesting questions a user might ask a chatbot. Return ONLY the questions, one per line:\n\n{pdf_text}"
    response = client.models.generate_content(model="gemini-2.5-flash", contents=prompt)
    raw_text = response.candidates[0].content.parts[0].text
    
    questions = [q.strip() for q in raw_text.split('\n') if q.strip()][:3]
    return {"suggested_questions": questions}

def rag_chat_node(state: GraphState):
    """Answers user queries by retrieving context from Pinecone."""
    print("---NODE: RAG CHAT---")
    user_query = state.get("user_query", "")
    session_id = state.get("session_id")

    vectorstore = get_vectorstore()
    search_kwargs = {"k": 3}
    if session_id:
        search_kwargs["filter"] = {"session_id": {"$eq": session_id}}
    docs = vectorstore.similarity_search(user_query, **search_kwargs)
    context = "\n".join([d.page_content for d in docs])
    
    prompt = f"""
    Context: {context}
    
    Question: {user_query}
    
    Answer to the question based on the context. the answer should be direct will very very precise elaboration(in mere 1 or 2 sentences) .
    You will answer questions which are related to the topic and surrounding knowledge use (if asked like elaborate more or give more details, not necessary within the context but surrounding question can be answered within the topic)
    you will use your llm knowlwdge to answer surrounding possible questions related to the topic and context.
    If the answer isn't there, or the user askes question which are no where related to the context or surrounding the topic say 'Your PDF does not have that content'.
    
    IN the end in seperate line suggest 2 more intresting questions or activities or anything based on the context based on the question and answer you provide to make the user more engaging.
    """
    response = client.models.generate_content(model="gemini-2.5-flash", contents=prompt)
    answer = response.candidates[0].content.parts[0].text
    
    return {"chat_response": answer}




# """Creates a structured outline for a PPT based on the summary text."""
def ppt_outline_node(state: GraphState):
    
    print("---NODE: CREATING PPT OUTLINE---")
    summary_text = state["summary_text"]
    
    prompt = f"""
    Create a 8-slide PowerPoint outline based on this summary content. 
    
    Return the response as a JSON array where each slide has:
    - "slide_number": number (1-8)
    - "title": slide title
    - "caption": brief description of the slide
    - "bullets": array of 3-4 bullet points
    
    Content to base slides on: {summary_text}
    
    Make it comprehensive and professional.
    """
    
    response = client.models.generate_content(
        model="gemini-2.5-flash", 
        contents=prompt,
        config={"response_mime_type": "application/json"}
    )
    
    ppt_outline = parse_json_safe(response.text)
    
    return {"ppt_outline": ppt_outline}




# """Creates the actual PPTX file from the outline data."""
def generate_ppt_file_node(state: GraphState):
    
    print("---NODE: GENERATING PPTX FILE---")
    from pptx import Presentation
    import uuid
    
    ppt_outline = state["ppt_outline"]
    
    if not ppt_outline:
        raise ValueError("No PPT outline found in state")
    
    
    prs = Presentation()
    
    #making slides based on outline
    for slide_data in ppt_outline:
        
        slide_layout = prs.slide_layouts[1]  # Title and Content 
        slide = prs.slides.add_slide(slide_layout)
        
       
        title = slide.shapes.title
        title.text = slide_data["title"]
        
        
        content = slide.placeholders[1]
        content.text = slide_data["caption"]
        
        
        text_frame = content.text_frame
        
        
        for bullet in slide_data["bullets"]:
            p = text_frame.add_paragraph()
            p.text = bullet
            p.level = 1
    
    # Save the presentation
    os.makedirs("assets", exist_ok=True)
    session_id = str(uuid.uuid4())[:8]
    ppt_path = f"assets/presentation_{session_id}.pptx"
    prs.save(ppt_path)
    
    return {"ppt_file_path": ppt_path}